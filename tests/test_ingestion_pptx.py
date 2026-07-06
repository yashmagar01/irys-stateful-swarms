"""Unit tests for pptx ingestion (S-12).

Builds a minimal PowerPoint deck with python-pptx, then verifies
that read_pptx() extracts slide text correctly.
"""
from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from src.ingestion.pptx import read_pptx


@pytest.fixture()
def simple_pptx(tmp_path: Path) -> Path:
    """Two-slide deck with known text on each slide."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # blank layout

    slide1 = prs.slides.add_slide(blank_layout)
    txBox1 = slide1.shapes.add_textbox(Inches(0), Inches(0), Inches(5), Inches(1))
    txBox1.text_frame.text = "Slide one content"

    slide2 = prs.slides.add_slide(blank_layout)
    txBox2 = slide2.shapes.add_textbox(Inches(0), Inches(0), Inches(5), Inches(1))
    txBox2.text_frame.text = "Slide two content"

    p = tmp_path / "deck.pptx"
    prs.save(str(p))
    return p


@pytest.fixture()
def empty_pptx(tmp_path: Path) -> Path:
    """A deck with one slide and no text."""
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    p = tmp_path / "empty.pptx"
    prs.save(str(p))
    return p


class TestReadPptxText:
    def test_slide_one_text_extracted(self, simple_pptx):
        text, _ = read_pptx(simple_pptx)
        assert "Slide one content" in text

    def test_slide_two_text_extracted(self, simple_pptx):
        text, _ = read_pptx(simple_pptx)
        assert "Slide two content" in text

    def test_returns_string(self, simple_pptx):
        text, meta = read_pptx(simple_pptx)
        assert isinstance(text, str)
        assert isinstance(meta, dict)

    def test_empty_deck_returns_string(self, empty_pptx):
        text, _ = read_pptx(empty_pptx)
        assert isinstance(text, str)


class TestReadPptxErrorHandling:
    def test_nonexistent_file_returns_error_string(self, tmp_path):
        bad = tmp_path / "missing.pptx"
        text, meta = read_pptx(bad)
        assert "error" in text.lower()
        assert meta == {}
