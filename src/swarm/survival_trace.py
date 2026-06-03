from __future__ import annotations

import json
from pathlib import Path

from .verification import extract_verification_targets, value_survives_in_text


def write_pending_survival_trace(
    output_dir: str,
    derived_report: dict | None,
    must_include: list[dict],
) -> None:
    """Write a pending trace after obligations/curation select final items."""
    if not output_dir or not derived_report:
        return
    items = []
    for item in derived_report.get("items", []):
        if not isinstance(item, dict):
            continue
        created_ids = item.get("created_entry_ids") or []
        obligation_ids = _matching_must_include_ids(created_ids, must_include)
        items.append({
            "derived_work_id": item.get("id"),
            "created_entry_ids": created_ids,
            "obligation_ids": obligation_ids,
            "obligated": bool(obligation_ids),
            "curated": bool(obligation_ids),
            "target_files": item.get("target_deliverables") or [],
            "found_in_artifact": False,
            "artifact_locations": [],
            "death_mode": _pending_death_mode(item, obligation_ids),
            "expected_values": _expected_values(item),
            "expected_text": item.get("reason", ""),
        })
    trace = {
        "schema_version": 1,
        "items": items,
        "summary": _summarize(items),
    }
    swarm_dir = Path(output_dir) / "swarm"
    swarm_dir.mkdir(parents=True, exist_ok=True)
    (swarm_dir / "commitment_survival_trace.pending.json").write_text(
        json.dumps(trace, indent=2),
        encoding="utf-8",
    )


def finalize_survival_trace(output_dir: str | Path,
                            artifact_texts: dict[str, str]) -> dict:
    """Finalize survival trace after deliverable files have been written."""
    out_dir = Path(output_dir)
    finalize_artifact_placement_trace(out_dir, artifact_texts)
    pending_path = out_dir / "swarm" / "commitment_survival_trace.pending.json"
    if not pending_path.exists():
        return {}
    trace = json.loads(pending_path.read_text(encoding="utf-8"))
    for item in trace.get("items", []):
        if not isinstance(item, dict):
            continue
        locations = _artifact_locations(item, artifact_texts)
        item["artifact_locations"] = locations
        item["found_in_artifact"] = bool(locations)
        if locations:
            item["death_mode"] = None
        elif not item.get("created_entry_ids"):
            item["death_mode"] = item.get("death_mode") or "executed_no_entry"
        elif not item.get("obligated"):
            item["death_mode"] = item.get("death_mode") or "not_obligated"
        else:
            item["death_mode"] = "artifact_missing"
    trace["summary"] = _summarize(trace.get("items", []))
    final_path = out_dir / "swarm" / "commitment_survival_trace.json"
    final_path.write_text(json.dumps(trace, indent=2), encoding="utf-8")
    return trace


def finalize_artifact_placement_trace(
    output_dir: str | Path,
    artifact_texts: dict[str, str],
) -> dict:
    """Trace artifact-native commitments after deliverable files are written."""
    out_dir = Path(output_dir)
    commitments_path = out_dir / "swarm" / "artifact_commitments.json"
    if not commitments_path.exists():
        return {}
    try:
        report = json.loads(commitments_path.read_text(encoding="utf-8"))
    except json.JSONDecodeError:
        return {}

    items = []
    for commitment in report.get("items", []):
        if not isinstance(commitment, dict):
            continue
        item = _artifact_commitment_trace_item(commitment, artifact_texts)
        items.append(item)

    trace = {
        "schema_version": 1,
        "items": items,
        "summary": _summarize_artifact_placements(items),
    }
    swarm_dir = out_dir / "swarm"
    swarm_dir.mkdir(parents=True, exist_ok=True)
    (swarm_dir / "artifact_placement_trace.json").write_text(
        json.dumps(trace, indent=2),
        encoding="utf-8",
    )
    return trace


def extract_artifact_texts(output_dir: str | Path,
                           deliverable_files: list[str]) -> dict[str, str]:
    output_path = Path(output_dir)
    texts: dict[str, str] = {}
    for filename in deliverable_files:
        path = output_path / filename
        if not path.exists():
            texts[filename] = ""
            continue
        suffix = path.suffix.lower()
        if suffix == ".docx":
            texts[filename] = _read_docx_text(path)
        elif suffix == ".xlsx":
            texts[filename] = _read_xlsx_text(path)
        elif suffix == ".pptx":
            texts[filename] = _read_pptx_text(path)
        else:
            texts[filename] = path.read_text(encoding="utf-8", errors="ignore")
    return texts


def _artifact_commitment_trace_item(
    commitment: dict,
    artifact_texts: dict[str, str],
) -> dict:
    target_file = str(commitment.get("target_file", "") or "")
    native_form = str(commitment.get("native_form", "") or "")
    verification_terms = _artifact_verification_terms(commitment)
    locations = _artifact_commitment_locations(
        verification_terms, commitment, artifact_texts,
    )
    target_locations = [
        loc for loc in locations
        if not target_file or loc.get("file") == target_file
    ]
    found_in_target = bool(target_locations)
    found_elsewhere = bool(locations) and not found_in_target
    wrong_format = _native_form_wrong_for_target(native_form, target_file)
    death_mode = None
    if wrong_format:
        death_mode = "wrong_format"
    elif found_in_target:
        death_mode = None
    elif found_elsewhere:
        death_mode = "wrong_file"
    elif target_file and target_file not in artifact_texts:
        death_mode = "artifact_missing"
    else:
        death_mode = "artifact_missing"

    return {
        "entry_id": commitment.get("entry_id", ""),
        "evidence_entry_ids": commitment.get("evidence_entry_ids", []),
        "target_file": target_file,
        "native_form": native_form,
        "artifact_function": commitment.get("artifact_function", ""),
        "satisfaction_conditions": commitment.get("satisfaction_conditions", []),
        "required_source_refs": commitment.get("required_source_refs", commitment.get("source_refs", [])),
        "verification_terms": verification_terms,
        "found_in_target_file": found_in_target,
        "found_elsewhere": found_elsewhere,
        "artifact_locations": locations,
        "death_mode": death_mode,
        "summary": commitment.get("summary", ""),
        "source": commitment.get("source", ""),
    }


def _artifact_commitment_locations(
    verification_terms: list[str],
    commitment: dict,
    artifact_texts: dict[str, str],
) -> list[dict]:
    if not verification_terms:
        return []
    context_terms = _context_terms({
        "expected_text": commitment.get("summary", ""),
        "required_inputs": [],
    })
    locations = []
    for filename, text in artifact_texts.items():
        for term in verification_terms:
            if value_survives_in_text(term, text, context_terms):
                locations.append({"file": filename, "evidence": term})
                break
    return locations


def _artifact_verification_terms(commitment: dict) -> list[str]:
    terms = []
    raw_terms = commitment.get("verification_terms", [])
    if isinstance(raw_terms, list):
        for term in raw_terms:
            clean = str(term).strip()
            if clean and clean not in terms:
                terms.append(clean)
    for target in extract_verification_targets(str(commitment.get("summary", ""))):
        raw = str(target.get("raw", "")).strip()
        if raw and raw not in terms:
            terms.append(raw)
    return terms[:12]


def _native_form_wrong_for_target(native_form: str, target_file: str) -> bool:
    if not native_form or not target_file:
        return False
    lower = target_file.lower()
    if native_form == "workbook_row":
        return not lower.endswith((".xlsx", ".csv"))
    if native_form == "slide_bullet":
        return not lower.endswith(".pptx")
    if native_form == "drafting_clause":
        return not any(word in lower for word in ("redline", "markup", "rider", ".docx"))
    return False


def _matching_must_include_ids(created_ids: list[str],
                               must_include: list[dict]) -> list[str]:
    created = set(created_ids)
    obligation_ids = []
    for idx, item in enumerate(must_include, 1):
        if not isinstance(item, dict):
            continue
        raw = item.get("entry_id") or item.get("entry_ids") or ""
        entry_ids = []
        if isinstance(raw, list):
            entry_ids = [str(v) for v in raw]
        else:
            entry_ids = [part.strip() for part in str(raw).split(",")]
        if created & set(entry_ids):
            obligation_ids.append(f"obl_{idx:03d}")
    return obligation_ids


def _pending_death_mode(item: dict, obligation_ids: list[str]) -> str | None:
    if item.get("status") in {"diagnostic_only", "executable"}:
        return "selected_but_not_executed"
    if item.get("status") == "execution_failed":
        return item.get("death_mode") or "executed_no_entry"
    if item.get("status") == "executed" and not obligation_ids:
        return "not_obligated"
    return None


def _expected_values(item: dict) -> list[str]:
    values = []
    for key in ("expected_result", "result", "expression"):
        raw = item.get(key)
        if isinstance(raw, str) and raw.strip():
            values.append(raw.strip())
    return values


def _artifact_locations(item: dict,
                        artifact_texts: dict[str, str]) -> list[dict]:
    values = [v for v in item.get("expected_values", []) if isinstance(v, str)]
    context_terms = _context_terms(item)
    locations = []
    for filename, text in artifact_texts.items():
        for value in values:
            if value_survives_in_text(value, text, context_terms):
                locations.append({"file": filename, "evidence": value})
                break
    return locations


def _context_terms(item: dict) -> list[str]:
    words = []
    for raw in [item.get("expected_text", "")] + [
        str(inp.get("label", ""))
        for inp in item.get("required_inputs", [])
        if isinstance(inp, dict)
    ]:
        for word in str(raw).split():
            cleaned = "".join(ch for ch in word.lower() if ch.isalpha())
            if len(cleaned) >= 4 and cleaned not in words:
                words.append(cleaned)
            if len(words) >= 6:
                return words
    return words


def _summarize(items: list[dict]) -> dict:
    selected = len(items)
    created = sum(1 for item in items if item.get("created_entry_ids"))
    obligated = sum(1 for item in items if item.get("obligated"))
    survived = sum(1 for item in items if item.get("found_in_artifact"))
    death_modes: dict[str, int] = {}
    for item in items:
        mode = item.get("death_mode")
        if mode:
            death_modes[mode] = death_modes.get(mode, 0) + 1
    return {
        "selected": selected,
        "created": created,
        "obligated": obligated,
        "artifact_survived": survived,
        "lost": selected - survived,
        "death_modes": death_modes,
    }


def _summarize_artifact_placements(items: list[dict]) -> dict:
    death_modes: dict[str, int] = {}
    native_forms: dict[str, int] = {}
    for item in items:
        mode = item.get("death_mode")
        if mode:
            death_modes[mode] = death_modes.get(mode, 0) + 1
        native = str(item.get("native_form") or "unknown")
        native_forms[native] = native_forms.get(native, 0) + 1
    found_target = sum(1 for item in items if item.get("found_in_target_file"))
    found_elsewhere = sum(1 for item in items if item.get("found_elsewhere"))
    return {
        "selected": len(items),
        "targeted": sum(1 for item in items if item.get("target_file")),
        "found_in_target_file": found_target,
        "found_elsewhere": found_elsewhere,
        "lost": len(items) - found_target,
        "death_modes": death_modes,
        "native_forms": native_forms,
    }


def _read_docx_text(path: Path) -> str:
    try:
        from docx import Document as DocxDocument
        doc = DocxDocument(str(path))
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception:
        return ""


def _read_xlsx_text(path: Path) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), data_only=False, read_only=True)
        parts = []
        for ws in wb.worksheets:
            parts.append(f"# Sheet: {ws.title}")
            for row in ws.iter_rows(values_only=True):
                values = ["" if value is None else str(value) for value in row]
                if any(values):
                    parts.append(" | ".join(values))
        return "\n".join(parts)
    except Exception:
        return ""


def _read_pptx_text(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts = []
        for idx, slide in enumerate(prs.slides, 1):
            parts.append(f"# Slide {idx}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
        return "\n".join(parts)
    except Exception:
        return ""
