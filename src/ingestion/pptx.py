from __future__ import annotations

from pathlib import Path


def read_pptx(path: Path) -> tuple[str, dict]:
    try:
        from markitdown import MarkItDown
        md = MarkItDown()
        result = md.convert(str(path))
        return result.text_content, {}
    except Exception:
        pass

    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"--- Slide {i} ---")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            parts.append(text)
        return "\n".join(parts), {}
    except Exception as e:
        return f"(error reading {path.name}: {e})", {}
